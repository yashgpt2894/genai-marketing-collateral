"""
Build the Case 2 presentation deck (<=10 slides, diagram-heavy) reflecting the
DEPLOYED system. Reproducible: python build_deck.py  ->  Case2_GenAI_Collateral.pptx

Design system (matches the project): teal #0FB5A6, ink #10151E, slate #475569,
Trebuchet MS (display) / Calibri (body) / Consolas (mono/eyebrow).
Needs python-pptx.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

TEAL, INK, INK2, SLATE, LIGHT, WHITE = "0FB5A6", "10151E", "1B2230", "475569", "F4F6F9", "FFFFFF"
ICE, LINE, MUTE, AMBER = "C7D0DE", "D7DEEA", "64748B", "E0A23E"
DISPLAY, BODY, MONO = "Trebuchet MS", "Calibri", "Consolas"
def rgb(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width, prs.slide_height = In(13.333), In(7.5)
BLANK = prs.slide_layouts[6]


def sld(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(bg)
    return s

def bar(s):  # teal left rail (dark slides)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(0), In(0), In(0.18), In(7.5))
    b.fill.solid(); b.fill.fore_color.rgb = rgb(TEAL); b.line.fill.background(); b.shadow.inherit = False

def _txt(s, text, x, y, w, h, *, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT,
         spacing=None, anchor=MSO_ANCHOR.TOP, caps=False):
    tb = s.shapes.add_textbox(In(x), In(y), In(w), In(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        r = p.add_run(); r.text = ln.upper() if caps else ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = rgb(color)
        if caps: r.font._rPr.set("spc", "260")
    return tb

def eyebrow(s, text, x=0.85, y=0.55, color=TEAL):
    _txt(s, text, x, y, 11, 0.3, size=12, color=color, bold=True, font=MONO, caps=True)

def title(s, text, x=0.82, y=0.95, w=11.6, color=INK, size=30):
    _txt(s, text, x, y, w, 1.1, size=size, color=color, bold=True, font=DISPLAY, spacing=1.02)

def footer(s, n, dark=False):
    _txt(s, "ML6 — Senior AI Engineer Challenge  ·  Case 2 · GenAI collateral",
         0.85, 7.08, 10, 0.3, size=9, color=(MUTE if not dark else "5A6B82"), font=MONO)
    _txt(s, str(n), 12.4, 7.08, 0.5, 0.3, size=9, color=(MUTE if not dark else "5A6B82"),
         font=MONO, align=PP_ALIGN.RIGHT)

def box(s, x, y, w, h, *, fill=WHITE, line=LINE, text="", tcolor=INK, size=12, bold=True,
        sub="", scolor=MUTE, font=BODY):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    if line: shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = In(0.08); tf.margin_right = In(0.08); tf.margin_top = In(0.04); tf.margin_bottom = In(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = rgb(tcolor)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(size - 3.5); r2.font.name = font; r2.font.color.rgb = rgb(scolor)
    return shp

def arrow(s, x1, y1, x2, y2, color=TEAL, width=1.75):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    c.line.color.rgb = rgb(color); c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c

def chip(s, x, y, w, text, color=TEAL, fill=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(x), In(y), In(w), In(0.34))
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    else: shp.fill.background()
    shp.line.color.rgb = rgb(color); shp.line.width = Pt(1); shp.shadow.inherit = False
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_top", "margin_bottom"): setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = rgb(color)
    return shp


# ── 1 · title ────────────────────────────────────────────────────────────────
s = sld(INK); bar(s)
eyebrow(s, "ML6 · SENIOR AI ENGINEER · CASE 2", 0.85, 1.5, TEAL)
title(s, "Automated Generative\nMarketing Collateral", 0.82, 2.05, 11.5, WHITE, 46)
_txt(s, "Company PDFs → on-brand, factually-grounded B2B articles, as structured JSON.",
     0.85, 4.35, 11, 0.5, size=17, color=ICE)
_txt(s, "A deterministic, grounded LLM pipeline — deployed end-to-end on Google Cloud.",
     0.85, 4.85, 11, 0.4, size=13, color="7E8CA0", font=MONO)
footer(s, 1, dark=True)

# ── 2 · the problem ───────────────────────────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "THE PROBLEM"); title(s, "Editors hand-build every personalized brochure")
steps = [("Research", "read both companies"), ("Write", "draft a bridging article"),
         ("Gather assets", "logos & images"), ("Fight the template", "fit columns, limits, colours")]
x = 0.85; w = 2.7; gap = 0.32
for i, (t, sub) in enumerate(steps):
    box(s, x, 2.5, w, 1.5, fill=LIGHT, line=LINE, text=t, sub=sub, size=14)
    if i < 3: arrow(s, x + w + 0.03, 3.25, x + w + gap - 0.03, 3.25, color=SLATE, width=1.5)
    x += w + gap
_txt(s, "Slow, repetitive, at volume — hours of editor time per brochure.\nAutomating research → layout turns that into seconds of compute at consistent, on-brand quality.",
     0.85, 4.5, 11.6, 1.0, size=15, color=SLATE, spacing=1.15)
footer(s, 2)

# ── 3 · solution = two endpoints + flow ───────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "THE SOLUTION · END-TO-END FLOW")
title(s, "Two endpoints, one deterministic flow")
# one continuous pipeline: input PDFs -> ... -> rendered brochure (the two POSTs marked along it)
flow = [
    ("Two PDFs", "sender + receiver", LINE),
    ("Parse", "PyMuPDF + Gemini", TEAL),
    ("Briefs", "cited facts", LINE),
    ("Generate", "draft → verify", TEAL),
    ("ArticleJSON", "layout-valid", LINE),
    ("Brochure", "render + review", LINE),
]
x = 0.85; w = 1.72; gap = 0.26; ytop = 3.05; h = 1.2; ymid = ytop + h / 2
xs = []
for i, (t, sub, ln) in enumerate(flow):
    xs.append(x)
    box(s, x, ytop, w, h, fill=LIGHT, line=ln, text=t, sub=sub, size=12)
    if i < len(flow) - 1:
        arrow(s, x + w + 0.01, ymid, x + w + gap - 0.01, ymid, color=SLATE, width=1.6)
    x += w + gap
# the two API endpoints, marked along the single flow
_txt(s, "POST /…/documents", xs[1] - 0.4, ytop - 0.46, w + 0.8, 0.3, size=10, color=TEAL, bold=True, font=MONO, align=PP_ALIGN.CENTER)
_txt(s, "POST /generate → 200", xs[3] - 0.5, ytop - 0.46, w + 1.0, 0.3, size=10, color=TEAL, bold=True, font=MONO, align=PP_ALIGN.CENTER)
_txt(s, "One continuous pipeline — input PDFs → grounded, cited ArticleJSON → rendered brochure. Deterministic, typed, fully traceable.",
     0.85, 4.55, 11.6, 0.6, size=14, color=SLATE, spacing=1.15)
chip(s, 0.85, 5.45, 3.0, "auth: Bearer ID token")
chip(s, 4.0, 5.45, 3.4, "structured JSON, schema-valid")
chip(s, 7.55, 5.45, 3.5, "every word limit enforced in code")
footer(s, 3)

# ── 4 · orchestration (the explicit requirement) ──────────────────────────────
s = sld(INK); bar(s)
eyebrow(s, "ORCHESTRATION", 0.85, 0.55, TEAL)
title(s, "Contextual retrieval → LLM interaction → layout", 0.82, 0.95, 11.6, WHITE)
stages = [
    ("1 · CONTEXTUAL RETRIEVAL", ["PDFs → distilled briefs", "source-tagged facts", "long-context (no RAG yet)"]),
    ("2 · LLM INTERACTION", ["draft — Gemini 2.5 Pro", "grounded, inline citations", "verify — 2.5 Flash"]),
    ("3 · LAYOUT / FORMATTING", ["map to template blocks", "validate limits in code", "repair loop → JSON"]),
]
x = 0.85; w = 3.7; gap = 0.45
for i, (hd, items) in enumerate(stages):
    box(s, x, 2.3, w, 3.0, fill=INK2, line="2C3647", text="", )
    _txt(s, hd, x + 0.25, 2.55, w - 0.5, 0.5, size=12, color=TEAL, bold=True, font=MONO)
    _txt(s, "\n".join("•  " + it for it in items), x + 0.25, 3.25, w - 0.5, 1.8, size=14, color=ICE, spacing=1.3)
    if i < 2: arrow(s, x + w + 0.05, 3.8, x + w + gap - 0.05, 3.8, color=TEAL, width=2)
    x += w + gap
_txt(s, "A fixed DAG — deterministic, typed, traceable. Not an open-ended agent.",
     0.85, 5.6, 11, 0.4, size=14, color="7E8CA0", font=MONO)
footer(s, 4, dark=True)

# ── 5 · cloud architecture ────────────────────────────────────────────────────
s = sld(INK); bar(s)
eyebrow(s, "ARCHITECTURE · GOOGLE CLOUD", 0.85, 0.55, TEAL)
title(s, "Managed, stateless, scales to zero", 0.82, 0.95, 11.6, WHITE)
box(s, 0.85, 2.15, 2.5, 0.8, fill=INK2, line="2C3647", text="Client / curl", sub="Bearer ID token", size=12, tcolor=WHITE, scolor=ICE)
box(s, 4.0, 2.15, 5.2, 0.8, fill="0B2E2A", line=TEAL, text="Cloud Run · API (FastAPI)", sub="auth · stateless · scale-to-zero", size=13, tcolor=WHITE, scolor=ICE)
arrow(s, 3.4, 2.55, 3.95, 2.55, color=TEAL)
# upload branch
box(s, 4.0, 3.4, 5.2, 0.85, fill=INK2, line="2C3647", text="Pub/Sub → parse worker", sub="PyMuPDF + Gemini 2.5 Flash", size=12, tcolor=WHITE, scolor=ICE)
arrow(s, 6.6, 2.95, 6.6, 3.38, color=TEAL)
# generate branch
box(s, 9.7, 2.15, 2.8, 0.85, fill=INK2, line="2C3647", text="Vertex AI", sub="Gemini 2.5 Pro / Flash", size=12, tcolor=WHITE, scolor=ICE)
arrow(s, 9.2, 2.55, 9.65, 2.55, color=TEAL)
# stores
box(s, 4.0, 4.7, 2.5, 0.85, fill=INK2, line="2C3647", text="Firestore", sub="briefs · jobs · outputs", size=12, tcolor=WHITE, scolor=ICE)
box(s, 6.7, 4.7, 2.5, 0.85, fill=INK2, line="2C3647", text="Cloud Storage", sub="PDFs · assets · CMEK", size=12, tcolor=WHITE, scolor=ICE)
arrow(s, 6.0, 4.25, 5.3, 4.68, color=SLATE, width=1.4)
arrow(s, 7.2, 4.25, 7.9, 4.68, color=SLATE, width=1.4)
# ── callout: visually highlight the core architectural decision (the API box) ──
cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(0.85), In(3.55), In(2.95), In(1.7))
cb.fill.solid(); cb.fill.fore_color.rgb = rgb(INK2)
cb.line.color.rgb = rgb(AMBER); cb.line.width = Pt(2); cb.shadow.inherit = False
_txt(s, "WHY CLOUD RUN", 1.08, 3.72, 2.6, 0.3, size=11, color=AMBER, bold=True, font=MONO, caps=True)
_txt(s, "Deterministic + tool-less → one stateless service. No agent runtime, no servers — scale-to-zero, tiny blast radius.",
     1.08, 4.08, 2.55, 1.1, size=11.5, color=ICE, spacing=1.16)
arrow(s, 2.45, 3.55, 4.05, 2.82, color=AMBER, width=2)
# callout: Vertex — no API key
vb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(9.6), In(3.12), In(2.9), In(0.74))
vb.fill.solid(); vb.fill.fore_color.rgb = rgb(INK2); vb.line.color.rgb = rgb(AMBER); vb.line.width = Pt(1.5); vb.shadow.inherit = False
_txt(s, "NO API KEY", 9.82, 3.22, 2.6, 0.28, size=10, color=AMBER, bold=True, font=MONO, caps=True)
_txt(s, "Vertex via the Cloud Run service account (ADC)", 9.82, 3.5, 2.5, 0.32, size=10, color=ICE, spacing=1.05)
# callout: async path
ab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(9.6), In(3.98), In(2.9), In(1.05))
ab.fill.solid(); ab.fill.fore_color.rgb = rgb(INK2); ab.line.color.rgb = rgb(AMBER); ab.line.width = Pt(1.5); ab.shadow.inherit = False
_txt(s, "ASYNC PARSE", 9.82, 4.08, 2.6, 0.28, size=10, color=AMBER, bold=True, font=MONO, caps=True)
_txt(s, "Pub/Sub queues parse/upload bursts (+ DLQ), off the request path.",
     9.82, 4.38, 2.55, 0.6, size=10, color=ICE, spacing=1.12)
arrow(s, 9.58, 4.35, 9.22, 3.92, color=AMBER, width=1.5)
# highlight: the security / governance spine
sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(0.78), In(5.82), In(11.75), In(0.86))
sp.fill.background(); sp.line.color.rgb = rgb(AMBER); sp.line.width = Pt(1.25); sp.shadow.inherit = False
_txt(s, "SECURITY SPINE", 1.0, 5.92, 3, 0.3, size=10, color=AMBER, bold=True, font=MONO, caps=True)
_txt(s, "Google ID-token auth + IAM   ·   CMEK   ·   EU-region storage   ·   Pub/Sub DLQ   ·   Terraform (infra/)",
     1.0, 6.26, 11.45, 0.5, size=11, color="9FB0C8", font=MONO, spacing=1.12)
footer(s, 5, dark=True)

# ── 6 · factual correctness ───────────────────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "THE CRUX · FACTUAL CORRECTNESS")
title(s, "Grounded, cited, verified — and honest")
fc = [("Briefs only", "distilled facts,\nnever raw PDF"), ("Inline citations", "every claim →\na brief fact-id"),
      ("Verify", "claim-level check,\nflag unsupported"), ("Human review", "reviewable draft;\nnever auto-publish")]
x = 0.85; w = 2.7; gap = 0.32
for i, (t, sub) in enumerate(fc):
    box(s, x, 2.5, w, 1.5, fill=LIGHT, text=t, sub=sub, size=14, scolor=SLATE)
    if i < 3: arrow(s, x + w + 0.03, 3.25, x + w + gap - 0.03, 3.25, color=SLATE, width=1.4)
    x += w + gap
_txt(s, "Grounding reduces hallucination — it doesn't eliminate it. So citations stay, and a human approves before anything is mailed.",
     0.85, 4.55, 11.6, 0.6, size=15, color=SLATE, spacing=1.15)
footer(s, 6)

# ── 7 · layout limits ─────────────────────────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "LAYOUT = HARD LIMITS")
title(s, "Don't trust the model to count")
box(s, 0.85, 2.6, 2.5, 1.1, fill=LIGHT, text="Draft block", sub="controlled generation", size=13)
arrow(s, 3.4, 3.15, 4.05, 3.15, color=SLATE, width=1.4)
box(s, 4.1, 2.6, 3.0, 1.1, fill="EAF8F6", line=TEAL, text="Validate in code", sub="words · image · colour", size=13, tcolor=INK)
arrow(s, 7.15, 3.15, 7.8, 3.15, color=SLATE, width=1.4)
box(s, 7.85, 2.6, 2.4, 1.1, fill=LIGHT, text="Repair (capped)", sub="rewrite offending block", size=13)
arrow(s, 9.05, 3.7, 9.05, 4.25, color=SLATE, width=1.4)
box(s, 7.85, 4.3, 2.4, 1.0, fill=LIGHT, text="Truncate", sub="sentence boundary", size=13)
# loop back
arrow(s, 7.84, 4.8, 5.6, 4.8, color=SLATE, width=1.4)
arrow(s, 5.6, 4.8, 5.6, 3.72, color=SLATE, width=1.4)
box(s, 10.5, 2.6, 1.9, 1.1, fill="0B2E2A", line=TEAL, text="valid JSON", size=13, tcolor=WHITE)
arrow(s, 10.25, 3.15, 10.45, 3.15, color=TEAL)
_txt(s, "Word/char limits, image placeholders and theme colours are ground-truth checks in a framework-free engine — fixed by a targeted loop, with graceful truncation as the floor.",
     0.85, 5.7, 11.6, 0.7, size=14, color=SLATE, spacing=1.15)
footer(s, 7)

# ── 8 · security & scale ──────────────────────────────────────────────────────
s = sld(INK); bar(s)
eyebrow(s, "SECURITY & SCALE", 0.85, 0.55, TEAL)
title(s, "Secure by default, scales on its own", 0.82, 0.95, 11.6, WHITE)
box(s, 0.85, 2.3, 5.55, 3.4, fill=INK2, line="2C3647")
_txt(s, "SECURITY", 1.15, 2.55, 5, 0.4, size=12, color=TEAL, bold=True, font=MONO)
_txt(s, "•  Google ID-token auth + Cloud Run IAM (defense in depth)\n•  Untrusted PDFs → injection hard gate (refuse, fail closed) + data-not-instructions channel\n•  CMEK-encrypted storage · EU region (europe-west1) · least-privilege IAM\n•  no tools / side-effects → tiny blast radius",
     1.15, 3.15, 5.0, 2.3, size=13.5, color=ICE, spacing=1.35)
box(s, 6.95, 2.3, 5.55, 3.4, fill=INK2, line="2C3647")
_txt(s, "SCALE", 7.25, 2.55, 5, 0.4, size=12, color=TEAL, bold=True, font=MONO)
_txt(s, "•  Cloud Run autoscales 0 → N on concurrency\n•  Pub/Sub absorbs upload bursts (queue + DLQ)\n•  long/bursty jobs → ?async=true (off-path)\n•  stateless → any instance serves any request\n•  ceiling = Vertex quota, not compute",
     7.25, 3.15, 5.0, 2.3, size=13.5, color=ICE, spacing=1.35)
footer(s, 8, dark=True)

# ── 9 · live demo (terminal) ──────────────────────────────────────────────────
s = sld(INK); bar(s)
eyebrow(s, "LIVE · END-TO-END ON GCP", 0.85, 0.55, TEAL)
title(s, "It runs — from the terminal", 0.82, 0.95, 11.6, WHITE)
term = ("$ curl -H \"Authorization: Bearer $TOKEN\" \\\n"
        "    -F role=sender -F files=@sender.pdf  $URL/companies/demo/documents\n"
        "$ curl ... -F role=receiver -F files=@receiver.pdf  .../documents\n"
        "$ curl -X POST $URL/generate -d '{\"pair_id\":\"demo\",\"prompt\":\"…\"}'\n"
        "    → 200  { template_id, blocks:[…], confidence, meta:{…} }")
b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, In(0.85), In(2.15), In(7.1), In(2.7))
b.fill.solid(); b.fill.fore_color.rgb = rgb("0A0E16"); b.line.color.rgb = rgb("2C3647"); b.shadow.inherit = False
_txt(s, term, 1.05, 2.35, 6.8, 2.4, size=11, color="9FE8DF", font=MONO, spacing=1.25)
# metrics callouts
metrics = [("gemini-2.5-pro", "writer model"), ("$0.01", "per brochure")]
y = 2.15
for t, sub in metrics:
    box(s, 8.35, y, 4.15, 0.72, fill=INK2, line="2C3647", text=t, sub=sub, size=15, tcolor=TEAL, scolor=ICE)
    y += 0.82
_txt(s, "Output, grounded in the PDFs:\n“Turn Vanguard's idle time into delivery time.”",
     0.85, 5.2, 7.1, 0.9, size=13.5, color=ICE, spacing=1.15)
footer(s, 9, dark=True)

# ── 10 · why these choices ─────────────────────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "DESIGN DECISIONS")
title(s, "Why these choices?")
_txt(s, "Pragmatism over shiny — the simplest thing that works, with the upgrade path named.",
     0.85, 1.72, 11.6, 0.4, size=14, color=SLATE)
decisions = [
    ("Deterministic typed pipeline — not an agent",
     "predictable, testable, traceable; ADK / Agent Engine only when we add tools or a HITL pause"),
    ("Long-context grounding — not RAG",
     "a bounded two-company corpus fits in context — full traceability, no vector DB; RAG when it grows"),
    ("Managed Gemini on Vertex — not self-hosted",
     "ship now at ~$0.01 / brochure with no GPUs to run; revisit fine-tuning only at real volume"),
    ("Limits & images enforced in code — not the model",
     "the LLM is great at language, weak at counting/choosing — code owns word limits, image slots, JSON"),
    ("Cloud Run — not GKE or Agent Engine",
     "stateless + tool-less → scale-to-zero, a tiny blast radius and minimal ops"),
]
y = 2.35; h = 0.8; gap = 0.1
for t, why in decisions:
    box(s, 0.85, y, 11.6, h, fill=LIGHT, line=LINE)
    _txt(s, t, 1.15, y + 0.12, 11.0, 0.32, size=14, color=INK, bold=True, font=DISPLAY)
    _txt(s, why, 1.15, y + 0.46, 11.0, 0.3, size=11.5, color=SLATE)
    y += h + gap
footer(s, 10)

# ── 11 · roadmap + trade-offs ──────────────────────────────────────────────────
s = sld(WHITE)
eyebrow(s, "ROADMAP & TRADE-OFFS")
title(s, "A working slice today — a clear path to scale")
cards = [
    ("Multi-tenancy", "tenant from the token; per-tenant CMEK + budgets",
     "single-tenant ships now — keys already threaded, so it's config not a refactor"),
    ("Org dashboard", "every generation → BigQuery → Looker (cost, groundedness, throughput)",
     "the BigQuery export is built — only the Looker wiring is left"),
    ("Go agentic → Agent Engine", "tools (live research, image sourcing) or a HITL pause",
     "agents add cost + opacity — not until there are real tools to call"),
    ("Harden", "Model Armor injection filter · tune the judge · golden eval set in CI",
     "model calls use the global endpoint today — EU residency is the open gap"),
]
xs = [0.85, 6.75]; ys = [2.35, 4.5]; w = 5.7; h = 2.0
for i, (t, d, trade) in enumerate(cards):
    x = xs[i % 2]; y = ys[i // 2]
    box(s, x, y, w, h, fill=LIGHT, line=LINE)
    _txt(s, t, x + 0.3, y + 0.2, w - 0.6, 0.4, size=15, color=INK, bold=True, font=DISPLAY)
    _txt(s, d, x + 0.3, y + 0.64, w - 0.6, 0.7, size=12.5, color=SLATE, spacing=1.12)
    _txt(s, "trade-off:  " + trade, x + 0.3, y + 1.42, w - 0.6, 0.5, size=11, color=AMBER, spacing=1.1)
footer(s, 11)

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Case2_GenAI_Collateral.pptx")
prs.save(out)
print("wrote", out, "·", len(prs.slides._sldIdLst), "slides")
